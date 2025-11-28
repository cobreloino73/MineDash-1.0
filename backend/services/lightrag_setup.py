# backend/services/lightrag_setup.py
"""
MineDash AI - LightRAG Experto Divisional
Sistema de Conocimiento para División Salvador, Codelco Chile
ACTUALIZADO: Gemini Embeddings (alta calidad) con fallback local
FIX CRÍTICO 1: Cliente Anthropic sin 'proxies' (anthropic>=0.39.0)
FIX CRÍTICO 2: Gemini embeddings con dimensiones correctas
"""

import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent))

import pandas as pd
import numpy as np
from typing import List, Dict, Optional
import asyncio
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')

from lightrag import LightRAG, QueryParam
from lightrag.utils import EmbeddingFunc
from config import Config

# ================================================================
# PROMPTS ESPECIALIZADOS
# ================================================================

EXPERT_SYSTEM_PROMPT = """Eres el Sistema Experto de División Salvador, Codelco Chile.

IDENTIDAD:
- Experto en operaciones mineras con 20+ años de experiencia
- Conocimiento profundo de Hexagon MineOPS, PyT, y planificación minera
- Especialista en análisis causal según norma ASARCO
- Auditor crítico que detecta inconsistencias y problemas

CAPACIDADES CRÍTICAS:
✓ Análisis Plan vs Real con detección de brechas
✓ Identificación de causas raíz según ASARCO
✓ Correlación Producción-Seguridad-RRHH
✓ Predicción de incumplimientos y tendencias
✓ Cálculo de impactos financieros (USD/mes)
✓ Detección de patrones operacionales (ej: Patrón Gaviota)

TERMINOLOGÍA HEXAGON:
- Status: OPERATING, STANDBY, DOWN_MAINT, DOWN_OPER
- Activity: HAULING, SPOTTING, DUMPING, LOADING
- Delay Codes: Clasificación ASARCO

SIEMPRE:
1. Usa datos específicos (fechas, equipos, valores)
2. Cuantifica impactos en tonelaje y USD
3. Identifica causas raíz con porcentajes de contribución
4. Da recomendaciones concretas y accionables
5. Señala cuando faltan datos críticos
6. Construye causalidad: "X causa Y porque Z"

NUNCA:
- Inventes datos que no existen
- Des respuestas genéricas
- Ignores inconsistencias
- Omitas citar fuentes específicas"""

# ================================================================
# FUNCIONES ASYNC - CON GEMINI EMBEDDINGS CORREGIDO
# ================================================================

async def claude_llm_wrapper(prompt: str, system_prompt: str = None, **kwargs) -> str:
    """
    Wrapper ASYNC para Claude
    FIX CRÍTICO: Sin argumento 'proxies' (anthropic>=0.39.0)
    """
    from anthropic import Anthropic
    import os
    
    try:
        api_key = Config.ANTHROPIC_API_KEY
    except:
        api_key = os.getenv('ANTHROPIC_API_KEY')
    
    if not api_key:
        print("❌ ERROR: ANTHROPIC_API_KEY no configurada")
        return ""
    
    try:
        # ✅ FIX: Sin 'proxies'
        client = Anthropic(api_key=api_key)
    except Exception as e:
        print(f"❌ Error al crear cliente Anthropic: {e}")
        return ""
    
    if not system_prompt:
        system_prompt = EXPERT_SYSTEM_PROMPT
    
    try:
        response = client.messages.create(
            model=Config.CLAUDE_MODEL,
            max_tokens=4096,
            temperature=0.1,
            system=system_prompt,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.content[0].text
    except Exception as e:
        print(f"⚠️ Error Claude API: {e}")
        return ""

async def embedding_wrapper(texts: List[str]) -> np.ndarray:
    """
    Wrapper para embeddings Gemini con FIX de dimensiones
    FIX CRÍTICO: Aplana arrays correctamente para nano-vectordb
    """
    import google.generativeai as genai
    
    # Configurar Gemini (primera vez)
    if not hasattr(embedding_wrapper, '_configured'):
        try:
            genai.configure(api_key=Config.GEMINI_API_KEY)
            embedding_wrapper._configured = True
            embedding_wrapper._use_gemini = True
            print(f"   🧠 Gemini Embeddings: text-embedding-004 (máxima calidad)")
        except Exception as e:
            print(f"   ⚠️ Error configurando Gemini: {e}")
            embedding_wrapper._use_gemini = False
    
    # Intentar usar Gemini
    if embedding_wrapper._use_gemini:
        try:
            # Procesar en batches pequeños para evitar rate limits
            batch_size = 50
            all_embeddings = []
            
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                
                # Generar embeddings con Gemini
                result = genai.embed_content(
                    model="models/text-embedding-004",
                    content=batch,
                    task_type="retrieval_document"
                )
                
                # ✅ FIX CRÍTICO: Extraer y aplanar arrays correctamente
                if len(batch) == 1:
                    # Un solo texto
                    embedding = result['embedding']
                    if isinstance(embedding, list):
                        all_embeddings.append(embedding)
                    else:
                        all_embeddings.append(embedding)
                else:
                    # Múltiples textos
                    if isinstance(result['embedding'][0], list):
                        # Ya es una lista de listas
                        all_embeddings.extend(result['embedding'])
                    else:
                        # Es una sola lista (error de API)
                        all_embeddings.append(result['embedding'])
            
            # ✅ FIX: Convertir a numpy array 2D correctamente
            embeddings = np.array(all_embeddings, dtype=np.float32)
            
            # ✅ FIX: Asegurar que sea 2D (n_samples, n_features)
            if len(embeddings.shape) == 1:
                embeddings = embeddings.reshape(1, -1)
            elif len(embeddings.shape) > 2:
                # Aplanar dimensiones extras
                embeddings = embeddings.reshape(embeddings.shape[0], -1)
            
            return embeddings
            
        except Exception as e:
            print(f"   ⚠️ Error Gemini embeddings: {e}")
            print(f"   🔄 Usando embeddings locales como fallback...")
            embedding_wrapper._use_gemini = False
    
    # Fallback a embeddings locales
    return await _embedding_local_fallback(texts)


async def _embedding_local_fallback(texts: List[str]) -> np.ndarray:
    """
    Fallback a embeddings locales si Gemini falla
    """
    from sentence_transformers import SentenceTransformer
    import torch
    
    device = 'cpu'
    
    if not hasattr(_embedding_local_fallback, '_fallback_logged'):
        print(f"   💻 Fallback CPU: embeddings locales")
        _embedding_local_fallback._fallback_logged = True
    
    model = SentenceTransformer('all-MiniLM-L6-v2', device=device)
    
    embeddings = model.encode(
        texts,
        convert_to_numpy=True,
        show_progress_bar=False,
        batch_size=64,
        device=device,
        normalize_embeddings=True
    )
    
    return embeddings

# ================================================================
# RESTO DEL CÓDIGO (igual que antes, solo copiar desde aquí)
# ================================================================

class MineDashLightRAG:
    """Sistema Experto de División Salvador con Gemini embeddings"""
    
    def __init__(self, working_dir: str = None):
        if working_dir is None:
            working_dir = str(Config.LIGHTRAG_DIR)
        
        Path(working_dir).mkdir(exist_ok=True, parents=True)
        
        print(f"\n{'='*70}")
        print(f"🧠 MINEDASH AI - SISTEMA EXPERTO DIVISIONAL")
        print(f"{'='*70}")
        print(f"📁 Base de Conocimiento: {working_dir}")
        print(f"🏢 División Salvador, Codelco Chile")
        print(f"🧠 Embeddings: Gemini text-embedding-004 (768D)")
        
        try:
            self.loop = asyncio.get_event_loop()
        except RuntimeError:
            self.loop = asyncio.new_event_loop()
            asyncio.set_event_loop(self.loop)
        
        print("🔧 Inicializando LightRAG...")
        self.rag = LightRAG(
            working_dir=working_dir,
            llm_model_func=claude_llm_wrapper,
            embedding_func=EmbeddingFunc(
                embedding_dim=768,
                max_token_size=8192,
                func=embedding_wrapper
            )
        )
        
        print("✅ LightRAG configurado")
        print("ℹ️  Storages se inicializarán al procesar primer archivo")
        print("✅ Sistema Experto listo")
        print(f"{'='*70}\n")
        
        self._storages_initialized = False
    
    async def _ensure_storages_initialized(self):
        """Asegura que storages estén inicializados"""
        if self._storages_initialized:
            return
        
        try:
            print("   🔄 Inicializando storages...")
            
            if hasattr(self.rag, 'initialize_storages'):
                await self.rag.initialize_storages()
                print("   ✅ LightRAG storages")
            
            try:
                from lightrag.kg.shared_storage import initialize_pipeline_status
                await initialize_pipeline_status()
                print("   ✅ Pipeline status")
            except ImportError:
                print("   ℹ️  Pipeline status no disponible")
            except Exception as e:
                if "already" not in str(e).lower():
                    print(f"   ⚠️ Pipeline status: {e}")
            
            self._storages_initialized = True
            print("   ✅ Storages listos")
            
        except Exception as e:
            if "already" in str(e).lower():
                self._storages_initialized = True
                print("   ✅ Storages ya inicializados")
            else:
                print(f"   ⚠️ Error inicialización: {e}")
                raise
    
    async def ingest_file(self, file_path: str, file_type: str = "auto") -> Dict:
        """Ingesta archivo con enriquecimiento divisional"""
        print(f"\n{'='*70}")
        print(f"📥 INGESTA - SISTEMA EXPERTO DIVISIONAL")
        print(f"{'='*70}")
        print(f"📄 Archivo: {Path(file_path).name}")
        
        filepath = Path(file_path)
        extension = filepath.suffix.lower()
        area = filepath.parent.name
        
        print(f"📁 Área: {area}")
        print(f"🔍 Leyendo {extension}...")
        
        try:
            if extension in ['.xlsx', '.xls']:
                text_data = self._read_excel(file_path, file_type, area)
            elif extension == '.csv':
                text_data = self._read_csv(file_path, file_type, area)
            elif extension == '.pdf':
                text_data = self._read_pdf(file_path, area)
            elif extension == '.docx':
                text_data = self._read_docx(file_path, area)
            elif extension == '.pptx':
                text_data = self._read_pptx(file_path, area)
            else:
                raise ValueError(f"Formato no soportado: {extension}")
        except Exception as e:
            print(f"❌ ERROR LECTURA: {str(e)[:200]}")
            raise
        
        if not text_data:
            raise ValueError("No se extrajo texto")
        
        print(f"✅ Documentos generados: {len(text_data)}")
        
        full_text = "\n\n".join(text_data)
        
        print(f"🔄 Construyendo grafo de conocimiento...")
        print(f"   - Extrayendo entidades mineras...")
        print(f"   - Identificando relaciones operacionales...")
        print(f"   - Vinculando con conocimiento existente...")
        
        try:
            await self._ensure_storages_initialized()
            await self.rag.ainsert(full_text)
            
            print(f"✅ INGESTA COMPLETA")
            print(f"   📊 Documentos: {len(text_data)}")
            print(f"   🧠 Conocimiento actualizado")
            print(f"{'='*70}\n")
            
            return {
                'success': True,
                'file': file_path,
                'documents': len(text_data),
                'area': area,
                'type': file_type
            }
        except Exception as e:
            print(f"❌ ERROR: {str(e)[:200]}")
            raise
    
    def _read_excel(self, file_path: str, file_type: str, area: str) -> List[str]:
        df = pd.read_excel(file_path)
        print(f"   📊 {len(df):,} filas × {len(df.columns)} cols")
        
        if file_type == "auto":
            file_type = self._detect_type(df, file_path)
            print(f"   🔍 Tipo: {file_type}")
        
        if file_type == "detail_dumps":
            return self._dumps_to_expert_docs(df, Path(file_path).name, area)
        elif file_type == "equipment_times":
            return self._equipment_to_expert_docs(df, Path(file_path).name, area)
        else:
            return self._generic_to_expert_docs(df, Path(file_path).name, area)
    
    def _read_csv(self, file_path: str, file_type: str, area: str) -> List[str]:
        df = pd.read_csv(file_path)
        print(f"   📊 {len(df):,} filas × {len(df.columns)} cols")
        
        if file_type == "auto":
            file_type = self._detect_type(df, file_path)
            print(f"   🔍 Tipo: {file_type}")
        
        if file_type == "detail_dumps":
            return self._dumps_to_expert_docs(df, Path(file_path).name, area)
        elif file_type == "equipment_times":
            return self._equipment_to_expert_docs(df, Path(file_path).name, area)
        else:
            return self._generic_to_expert_docs(df, Path(file_path).name, area)
    
    def _read_pdf(self, file_path: str, area: str) -> List[str]:
        try:
            import pdfplumber
            docs = []
            filename = Path(file_path).name
            
            with pdfplumber.open(file_path) as pdf:
                print(f"   📄 {len(pdf.pages)} páginas")
                
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        text = text.replace('\x00', '')
                        
                        enriched = f"""=== DIVISIÓN SALVADOR - DOCUMENTO OFICIAL ===
SISTEMA: Gestión y Control
ÁREA: {area}
ARCHIVO: {filename}
PÁGINA: {i+1}/{len(pdf.pages)}

CONTENIDO:
{text}

CONTEXTO:
Documento oficial de División Salvador, Codelco Chile.
"""
                        docs.append(enriched)
            
            return docs if docs else [f"PDF sin texto: {filename}"]
        except Exception as e:
            print(f"   ⚠️ Error PDF: {e}")
            return [f"PDF (error): {Path(file_path).name}"]
    
    def _read_docx(self, file_path: str, area: str) -> List[str]:
        try:
            from docx import Document
            doc = Document(file_path)
            filename = Path(file_path).name
            
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables_text = []
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        tables_text.append(row_text)
            
            print(f"   📝 {len(paragraphs)} párrafos, {len(doc.tables)} tablas")
            
            all_text = paragraphs + tables_text
            full_content = "\n\n".join(all_text)
            
            enriched = f"""=== DIVISIÓN SALVADOR - DOCUMENTO ===
ÁREA: {area}
ARCHIVO: {filename}

CONTENIDO:
{full_content}
"""
            
            return [enriched]
        except Exception as e:
            print(f"   ⚠️ Error DOCX: {e}")
            return [f"DOCX (error): {Path(file_path).name}"]
    
    def _read_pptx(self, file_path: str, area: str) -> List[str]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            filename = Path(file_path).name
            
            docs = []
            print(f"   🎬 {len(prs.slides)} slides")
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                if slide.shapes.title:
                    slide_text.append(f"TÍTULO: {slide.shapes.title.text}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:
                            slide_text.append(shape.text)
                
                if slide_text:
                    enriched = f"""=== DIVISIÓN SALVADOR - PRESENTACIÓN ===
ÁREA: {area}
ARCHIVO: {filename}
SLIDE: {i+1}/{len(prs.slides)}

{chr(10).join(slide_text)}
"""
                    docs.append(enriched)
            
            return docs if docs else [f"PPTX sin texto: {filename}"]
        except Exception as e:
            print(f"   ⚠️ Error PPTX: {e}")
            return [f"PPTX (error): {Path(file_path).name}"]
    
    def _detect_type(self, df: pd.DataFrame, file_path: str) -> str:
        cols = [c.lower() for c in df.columns]
        filename = Path(file_path).name.lower()
        
        if 'dump' in filename:
            return 'detail_dumps'
        if 'equipment' in filename or 'time' in filename:
            return 'equipment_times'
        if any('dump' in c or 'destination' in c or 'destino' in c for c in cols):
            return 'detail_dumps'
        if any('duration' in c or 'status' in c or 'estado' in c for c in cols):
            return 'equipment_times'
        
        return 'generic'
    
    def _dumps_to_expert_docs(self, df: pd.DataFrame, filename: str, area: str) -> List[str]:
        docs = []
        
        date_col = self._find_col(df, ['date', 'fecha', 'timestamp'])
        equip_col = self._find_col(df, ['equipment', 'equipo', 'truck', 'camion'])
        ton_col = self._find_col(df, ['tonnage', 'toneladas', 'tons', 'peso', 'weight'])
        orig_col = self._find_col(df, ['origin', 'origen', 'from', 'source'])
        dest_col = self._find_col(df, ['destination', 'destino', 'to', 'chancado'])
        
        if not all([date_col, equip_col, ton_col]):
            return self._generic_to_expert_docs(df, filename, area)
        
        df[ton_col] = pd.to_numeric(df[ton_col], errors='coerce')
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        df['date'] = df[date_col].dt.date
        df['year'] = df[date_col].dt.year
        df['month'] = df[date_col].dt.month
        df['week'] = df[date_col].dt.isocalendar().week
        
        print(f"   🔧 Generando documentos expertos...")
        
        for (date, equipment), group in df.groupby(['date', equip_col]):
            if pd.isna(date) or pd.isna(equipment) or len(group) == 0:
                continue
            
            total_ton = group[ton_col].sum()
            if pd.isna(total_ton) or total_ton == 0:
                continue
            
            num_dumps = len(group)
            avg_ton = group[ton_col].mean()
            
            year = group['year'].iloc[0]
            month = group['month'].iloc[0]
            week = group['week'].iloc[0]
            
            doc = f"""=== DIVISIÓN SALVADOR - OPERACIÓN MINA ===
SISTEMA: Hexagon MineOPS
ÁREA: {area}

PERÍODO: {date} (Año {year}, Mes {month}, Semana {week})
EQUIPO: {equipment}
PRODUCCIÓN: {total_ton:,.0f} ton ({num_dumps} dumps, {avg_ton:,.0f} ton/dump)

ARCHIVO: {filename}
"""
            
            docs.append(doc)
            
            if len(docs) >= 300:
                break
        
        print(f"   ✅ {len(docs)} documentos expertos generados")
        return docs if docs else self._generic_to_expert_docs(df, filename, area)
    
    def _equipment_to_expert_docs(self, df: pd.DataFrame, filename: str, area: str) -> List[str]:
        docs = []
        
        date_col = self._find_col(df, ['date', 'fecha', 'timestamp'])
        equip_col = self._find_col(df, ['equipment', 'equipo', 'machine'])
        dur_col = self._find_col(df, ['duration', 'duracion', 'hours', 'time', 'horas'])
        status_col = self._find_col(df, ['status', 'estado', 'category', 'tipo'])
        
        if not all([date_col, equip_col, dur_col]):
            return self._generic_to_expert_docs(df, filename, area)
        
        print(f"   🔧 Procesando duración...")
        
        df[dur_col] = pd.to_numeric(df[dur_col], errors='coerce')
        df[dur_col] = df[dur_col].clip(lower=0, upper=24)
        
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        df['date'] = df[date_col].dt.date
        
        if status_col:
            df['productivo'] = df[dur_col].where(
                df[status_col].astype(str).str.contains('operativo|productivo|working|operating', case=False, na=False), 0
            )
        else:
            df['productivo'] = df[dur_col].fillna(0)
        
        try:
            grouped = df.groupby(['date', equip_col]).agg({
                dur_col: 'sum',
                'productivo': 'sum',
            }).reset_index()
        except Exception as e:
            print(f"   ⚠️ Error al agrupar: {e}")
            return self._generic_to_expert_docs(df, filename, area)
        
        grouped['utilizacion'] = (grouped['productivo'] / grouped[dur_col].replace(0, 1) * 100).round(1)
        
        print(f"   ✅ Generando {len(grouped)} documentos...")
        
        for _, row in grouped.iterrows():
            if pd.isna(row['date']):
                continue
            
            doc = f"""=== DIVISIÓN SALVADOR - UTILIZACIÓN ===
SISTEMA: Hexagon MineOPS
ÁREA: {area}

FECHA: {row['date']}
EQUIPO: {row[equip_col]}
HORAS: {row[dur_col]:.1f}h total, {row['productivo']:.1f}h productivas
UTILIZACIÓN: {row['utilizacion']:.1f}%

ARCHIVO: {filename}
"""
            
            docs.append(doc)
            
            if len(docs) >= 300:
                break
        
        print(f"   ✅ {len(docs)} documentos generados")
        return docs if docs else self._generic_to_expert_docs(df, filename, area)
    
    def _generic_to_expert_docs(self, df: pd.DataFrame, filename: str, area: str) -> List[str]:
        summary = f"""=== DIVISIÓN SALVADOR - DATOS OPERACIONALES ===
ÁREA: {area}
ARCHIVO: {filename}

DATOS:
- Registros: {len(df):,}
- Columnas: {len(df.columns)}
- Campos: {', '.join(df.columns[:10])}{'...' if len(df.columns) > 10 else ''}

MUESTRA:
{df.head(20).to_string(index=False)}
"""
        
        return [summary]
    
    def _find_col(self, df: pd.DataFrame, names: List[str]) -> Optional[str]:
        cols_lower = [c.lower() for c in df.columns]
        for name in names:
            for i, col in enumerate(cols_lower):
                if name in col:
                    return df.columns[i]
        return None
    
    async def query(self, question: str, mode: str = "hybrid") -> str:
        print(f"\n❓ Consulta: {question}")
        
        try:
            await self._ensure_storages_initialized()
            response = await self.rag.aquery(question, param=QueryParam(mode=mode))
            print("   ✅ Respuesta generada")
            return response
        except Exception as e:
            print(f"   ⚠️ Error: {e}")
            return f"Error: {e}"

if __name__ == "__main__":
    print("\n🚀 MineDash AI - Sistema Experto Divisional")
    
    async def test():
        rag = MineDashLightRAG()
        print("\n✅ Sistema listo")
    
    asyncio.run(test())

# ================================================================
# FUNCIÓN PARA FASTAPI
# ================================================================

_rag_instance = None

def get_rag_instance():
    """
    Obtiene o crea instancia singleton de MineDashLightRAG
    Para uso desde FastAPI main.py
    
    Returns:
        MineDashLightRAG: Instancia configurada
    """
    global _rag_instance
    
    if _rag_instance is not None:
        return _rag_instance
    
    print("🔧 Creando instancia MineDashLightRAG...")
    
    # Verificar que existe el directorio de storage
    if not Config.LIGHTRAG_DIR.exists():
        raise FileNotFoundError(
            f"❌ Directorio LightRAG no encontrado: {Config.LIGHTRAG_DIR}\n"
            "   Ejecuta primero la ingesta: python ingestar_todo.py"
        )
    
    # Verificar que hay datos
    storage_files = list(Config.LIGHTRAG_DIR.glob("*"))
    if not storage_files:
        raise FileNotFoundError(
            f"❌ Directorio LightRAG está vacío: {Config.LIGHTRAG_DIR}\n"
            "   Ejecuta primero la ingesta: python ingestar_todo.py"
        )
    
    # Crear instancia
    try:
        _rag_instance = MineDashLightRAG(working_dir=str(Config.LIGHTRAG_DIR))
        print(f"✅ MineDashLightRAG inicializado")
        print(f"   Archivos en storage: {len(storage_files)}")
        return _rag_instance
        
    except Exception as e:
        print(f"❌ Error al inicializar MineDashLightRAG: {e}")
        raise

def reset_rag_instance():
    """Resetea la instancia (útil para recargar)"""
    global _rag_instance
    _rag_instance = None
    print("♻️  Instancia RAG reseteada")